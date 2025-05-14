import os

import pytest
from presentation import Presentation
from pptx import Presentation as Pptx

@pytest.fixture
def pres():
    # Each test gets a fresh Presentation instance
    pres = Presentation()
    pres.add_slide(layout_index=6)
    return pres

def test_presentation_object_exists(pres):
    # Under the hood we should have a python-pptx Presentation
    assert pres._presentation is not None

def test_initial_slide_count(pres):
    # The very first slide‐deck should contain exactly 1 slide
    info = pres.get_presentation_info()
    assert "slide_count" in info
    assert info["slide_count"] == 1


def test_add_textbox_and_inspect_shape_info(pres):
    # add a textbox to slide 0
    r = pres.add_textbox(0, left=1, top=1, width=2, height=1, text="🐱 Hello!")
    assert "shape_index" in r
    info = pres.get_slide_info(0)
    # must find at least one TEXT_BOX in the shapes list
    assert any(s["shape_type"].startswith("TEXT_BOX") for s in info["shapes"])

def test_save_and_reload(tmp_path, pres):
    out = tmp_path / "demo.pptx"
    r = pres.save_presentation(str(out))
    assert os.path.isfile(str(out))
    # reload via python-pptx
    loaded = Pptx(str(out))
    assert len(loaded.slides) == pres.get_presentation_info()["slide_count"]

def test_add_and_move_shape(pres):
    # 1) Add a rectangle on slide 0
    resp = pres.add_shape(
        0,            # slide_index
        "rectangle",  # shape_type
        0.5,          # left
        0.5,          # top
        1.0,          # width
        1.0           # height
    )
    # should return a dict with at least "shape_index"
    assert isinstance(resp, dict)
    assert "shape_index" in resp
    shape_idx = resp["shape_index"]

    # 2) Now move that shape to a new position
    move_resp = pres.move_element(
        shape_idx,
        left=2.0,
        top=3.0
    )
    # move_element returns the new coordinates
    assert isinstance(move_resp, dict)
    assert move_resp["new_left"] == pytest.approx(2.0)
    assert move_resp["new_top"]  == pytest.approx(3.0)
