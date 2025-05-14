import os
import os
import subprocess
import tempfile
import uuid
from typing import Dict, List, Optional, Any

# — now it’s safe to import Aspose —
import aspose.slides as slides
from aspose.slides.export import SVGOptions, SvgExternalFontsHandling
from lxml import etree

import ppt_utils


def validate_parameters(params):
    """
    Validate parameters against constraints.
    
    Args:
        params: Dictionary of parameter name: (value, constraints) pairs
        
    Returns:
        (True, None) if all valid, or (False, error_message) if invalid
    """
    for param_name, (value, constraints) in params.items():
        for constraint_func, error_msg in constraints:
            if not constraint_func(value):
                return False, f"Parameter '{param_name}': {error_msg}"
    return True, None


def is_positive(value):
    """Check if a value is positive."""
    return value > 0


def is_non_negative(value):
    """Check if a value is non-negative."""
    return value >= 0


def is_in_range(min_val, max_val):
    """Create a function that checks if a value is in a range."""
    return lambda x: min_val <= x <= max_val


def is_in_list(valid_list):
    """Create a function that checks if a value is in a list."""
    return lambda x: x in valid_list


def is_valid_rgb(color_list):
    """Check if a color list is a valid RGB tuple."""
    if not isinstance(color_list, list) or len(color_list) != 3:
        return False
    return all(isinstance(c, int) and 0 <= c <= 255 for c in color_list)


def add_shape_direct(slide, shape_type: str, left: float, top: float, width: float, height: float) -> Any:
    """
    Add an auto shape to a slide using direct integer values instead of enum objects.
    
    This implementation provides a reliable alternative that bypasses potential 
    enum-related issues in the python-pptx library.
    
    Args:
        slide: The slide object
        shape_type: Shape type string (e.g., 'rectangle', 'oval', 'triangle')
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        The created shape
    """
    from pptx.util import Inches

    # Direct mapping of shape types to their integer values
    # These values are directly from the MS Office VBA documentation
    shape_type_map = {
        'rectangle': 1,
        'rounded_rectangle': 2,
        'oval': 9,
        'diamond': 4,
        'triangle': 5,  # This is ISOSCELES_TRIANGLE
        'right_triangle': 6,
        'pentagon': 56,
        'hexagon': 10,
        'heptagon': 11,
        'octagon': 12,
        'star': 12,  # This is STAR_5_POINTS (value 12)
        'arrow': 13,
        'cloud': 35,
        'heart': 21,
        'lightning_bolt': 22,
        'sun': 23,
        'moon': 24,
        'smiley_face': 17,
        'no_symbol': 19,
        'flowchart_process': 112,
        'flowchart_decision': 114,
        'flowchart_data': 115,
        'flowchart_document': 119
    }

    # Check if shape type is valid before trying to use it
    shape_type_lower = str(shape_type).lower()
    if shape_type_lower not in shape_type_map:
        available_shapes = ', '.join(sorted(shape_type_map.keys()))
        raise ValueError(f"Unsupported shape type: '{shape_type}'. Available shape types: {available_shapes}")

    # Get the integer value for the shape type
    shape_value = shape_type_map[shape_type_lower]

    # Create the shape using the direct integer value
    try:
        # The integer value is passed directly to add_shape
        shape = slide.shapes.add_shape(
            shape_value, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        return shape
    except Exception as e:
        raise ValueError(f"Failed to create '{shape_type}' shape using direct value {shape_value}: {str(e)}")


def remove_aspose_watermark(svg_xml: str) -> str:
    """Strip any <text>…Aspose…</text> via XML parsing (lxml)."""
    parser = etree.XMLParser(ns_clean=True, recover=True)
    tree = etree.fromstring(svg_xml.encode("utf-8"), parser=parser)
    ns = tree.nsmap.copy()
    if None in ns:
        ns["svg"] = ns.pop(None)
    for txt in tree.xpath(
            '//svg:text[contains(normalize-space(string(.)), "Aspose")]',
            namespaces=ns
    ):
        parent = txt.getparent()
        if parent is not None:
            parent.remove(txt)

    # re-add prolog + DOCTYPE so Inkscape parses it correctly
    decl = '<?xml version="1.0" encoding="utf-8" standalone="yes"?>\n'
    doctype = (
        '<!DOCTYPE svg PUBLIC "-//W3C//DTD SVG 1.1//EN" '
        '"http://www.w3.org/Graphics/SVG/1.1/DTD/svg11.dtd">\n'
    )
    body = etree.tostring(tree, encoding="utf-8").decode("utf-8")
    return decl + doctype + body

def svg_to_png(svg_xml: str, dpi: int = 300) -> bytes:
    """
    Save the cleaned SVG to a temp file, call Inkscape CLI to render
    it at `dpi`, read back the PNG bytes, and clean up.
    """
    # 1) write SVG out
    with tempfile.NamedTemporaryFile(suffix=".svg", delete=False) as svg_tmp:
        svg_tmp.write(svg_xml.encode("utf-8"))
        svg_path = svg_tmp.name

    # 2) designate a PNG temp
    png_fd, png_path = tempfile.mkstemp(suffix=".png")
    os.close(png_fd)

    # 3) render with Inkscape
    subprocess.run([
        "inkscape",
        svg_path,
        "--export-type=png",
        "--export-filename", png_path,
        "--export-dpi", str(dpi)
    ], check=True)

    # 4) read & cleanup
    png_bytes = open(png_path, "rb").read()
    os.remove(svg_path)
    os.remove(png_path)
    return png_bytes


# ---- Presentation Tools ----
class Presentation:
    def __init__(self):
        self.id = str(uuid.uuid4())
        self.slide = None
        # self._presentation = None

        # Create the session presentation
        # self.create_presentation()
        """Create a new PPTX and add one slide of the chosen layout."""
        self._presentation = ppt_utils.create_presentation()


    def get_layouts(self):
        layouts = [
            {"index": i, "name": layout.name}
            for i, layout in enumerate(self._presentation.slide_layouts)
        ]
        return {"layouts": layouts}

    def add_layout_index(self, layout_index: Optional[int] = None) -> Dict:
        # if caller passed a layout_index, use it; otherwise find the blank one
        if layout_index is None:
            # find the blank layout by name
            for i, layout in enumerate(self._presentation.slide_layouts):
                if layout.name.strip().lower() == "blank":
                    layout_index = i
                    break
            else:
                layout_index = 0  # fallback if no “Blank” found

        self.slide, _ = ppt_utils.add_slide(self._presentation, layout_index=layout_index)
        return {
            "presentation_id": self.id,
            "message": f"Created new presentation with layout {layout_index}: "
                       f"{self._presentation.slide_layouts[layout_index].name}",
            "slide_count": len(self._presentation.slides)
        }


    def _validate(
            self,
            *,
            layout_index: int | None = None,
            slide_index: int | None = None
    ) -> Dict | None:
        """
        Returns an error dict if:
         - no presentation is loaded
         - layout_index is out of range
         - slide_index is out of range
        Otherwise returns None.
        """
        if self._presentation is None:
            return {"error": "No presentation is currently loaded"}

        if layout_index is not None:
            max_layout = len(self._presentation.slide_layouts) - 1
            if not (0 <= layout_index <= max_layout):
                return {
                    "error": f"Invalid layout index: {layout_index}. "
                             f"Available layouts: 0–{max_layout}",
                    "available_layouts": ppt_utils.get_slide_layouts(self._presentation)
                }

        if slide_index is not None:
            max_slide = len(self._presentation.slides) - 1
            if not (0 <= slide_index <= max_slide):
                return {
                    "error": f"Invalid slide index: {slide_index}. "
                             f"Available slides: 0–{max_slide}"
                }

        return None

    def save_presentation(self, file_path: Optional[str] = None) -> Dict[str, Any]:
        """
        Save the PPTX to disk.
        If no file_path is provided, uses '{presentation_id}.pptx'.
        """
        # 1) validate
        if err := self._validate():
            return err

        # 2) default to <id>.pptx if no path supplied
        target = file_path or f"{self.id}.pptx"

        # 3) save and return a consistent payload
        try:
            saved = ppt_utils.save_presentation(self._presentation, target)
            return {
                "message": f"Saved to {saved}",
                "file_path": saved
            }
        except Exception as e:
            return {
                "error": f"Save failed: {e}"
            }

    def get_presentation_info(self) -> Dict:
        """Get information about a presentation."""

        # Validate presentation
        err = self._validate()
        if err:
            return err

        # pres = presentations[pres_id]

        # Get slide layouts
        layouts = ppt_utils.get_slide_layouts(self._presentation)

        # Get core properties
        core_props = ppt_utils.get_core_properties(self._presentation)

        return {
            "presentation_id": self.id,
            "slide_count": len(self._presentation.slides),
            "slide_layouts": layouts,
            "core_properties": core_props
        }

    # ---- Slide Tools ----
    def add_slide(
            self,
            layout_index: int = 1,
            title: Optional[str] = None
    ) -> Dict:
        # 1) validate layout
        if err := self._validate(layout_index=layout_index):
            return err

        # 2) do the add_slide
        slide, error = ppt_utils.safe_operation(
            "add_slide",
            lambda: ppt_utils.add_slide(self._presentation, layout_index)
        )
        if error:
            return {"error": error}

        # 3) set title if any...
        if title and slide[0].shapes.title:
            _, err2 = ppt_utils.safe_operation(
                "set_title",
                lambda: ppt_utils.set_title(slide[0], title)
            )
            if err2:
                return {
                    "warning": f"Slide created but failed to set title: {err2}",
                    "slide_index": len(self._presentation.slides) - 1,
                    "layout_name": slide[1].name
                }

        # 4) placeholders
        placeholders, err3 = ppt_utils.safe_operation(
            "get_placeholders",
            lambda: ppt_utils.get_placeholders(slide[0])
        )
        if err3:
            placeholders = []

        return {
            "message": f"Added slide with layout '{slide[1].name}'",
            "slide_index": len(self._presentation.slides) - 1,
            "layout_name": slide[1].name,
            "placeholders": placeholders
        }

    def get_slide_info(self, slide_index: int) -> Dict:
        # 1) validate slide index
        if err := self._validate(slide_index=slide_index):
            return err

        slide = self._presentation.slides[slide_index]
        placeholders = ppt_utils.get_placeholders(slide)

        shapes_info = [
            {
                "index": i,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "width": shape.width.inches,
                "height": shape.height.inches,
                "left": shape.left.inches,
                "top": shape.top.inches
            }
            for i, shape in enumerate(slide.shapes)
        ]

        return {
            "slide_index": slide_index,
            "placeholders": placeholders,
            "shapes": shapes_info
        }

    def populate_placeholder(
            self,
            slide_index: int,
            placeholder_idx: int,
            text: str,
    ) -> Dict:
        """Populate a placeholder with text."""
        # Validate presentation
        err = self._validate()
        if err:
            return err

        pres = self._presentation

        # Check if slide index is valid
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            # Check if placeholder exists
            if placeholder_idx not in [p.placeholder_format.idx for p in slide.placeholders]:
                return {
                    "error": f"Placeholder with index {placeholder_idx} not found in slide {slide_index}"
                }

            # Populate the placeholder
            ppt_utils.populate_placeholder(slide, placeholder_idx, text)

            return {
                "message": f"Populated placeholder {placeholder_idx} in slide {slide_index}"
            }
        except Exception as e:
            return {
                "error": f"Failed to populate placeholder: {str(e)}"
            }

    def add_bullet_points(
            self,
            slide_index: int,
            placeholder_idx: int,
            bullet_points: List[str],
    ) -> Dict:
        """Add bullet points to a placeholder."""

        # Validate presentation
        err = self._validate()
        if err:
            return err

        pres = self._presentation

        # Check if slide index is valid
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            # Check if placeholder exists
            if placeholder_idx not in [p.placeholder_format.idx for p in slide.placeholders]:
                return {
                    "error": f"Placeholder with index {placeholder_idx} not found in slide {slide_index}"
                }

            # Get the placeholder
            placeholder = slide.placeholders[placeholder_idx]

            # Add bullet points
            ppt_utils.add_bullet_points(placeholder, bullet_points)

            return {
                "message": f"Added {len(bullet_points)} bullet points to placeholder {placeholder_idx} in slide {slide_index}"
            }
        except Exception as e:
            return {
                "error": f"Failed to add bullet points: {str(e)}"
            }

    # ---- Text Tools ----
    def add_textbox(
            self,
            slide_index: int,
            left: float,
            top: float,
            width: float,
            height: float,
            text: str,
            font_size: Optional[int] = None,
            font_name: Optional[str] = None,
            bold: Optional[bool] = None,
            italic: Optional[bool] = None,
            color: Optional[List[int]] = None,
            alignment: Optional[str] = None,
    ) -> Dict:
        """Add a textbox to a slide."""
        # Validate presentation
        err = self._validate()
        if err:
            return err

        pres = self._presentation

        # Check if slide index is valid
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            # Add the textbox
            textbox = ppt_utils.add_textbox(slide, left, top, width, height, text)

            # Format the text if formatting options are provided
            if any([font_size, font_name, bold, italic, color, alignment]):
                ppt_utils.format_text(
                    textbox.text_frame,
                    font_size=font_size,
                    font_name=font_name,
                    bold=bold,
                    italic=italic,
                    color=tuple(color) if color else None,
                    alignment=alignment
                )

            return {
                "message": f"Added textbox to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1
            }
        except Exception as e:
            return {
                "error": f"Failed to add textbox: {str(e)}"
            }

    # ---- Image Tools ----
    def add_image(
            self,
            slide_index: int,
            image_path: str,
            left: float,
            top: float,
            width: Optional[float] = None,
            height: Optional[float] = None,
    ) -> Dict:
        """Add an image to a slide with graceful error recovery."""

        # Validate presentation
        err = self._validate()
        if err:
            return err

        pres = self._presentation

        # Check if slide index is valid
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        # Check if image file exists
        if not os.path.exists(image_path):
            # Try to find the image by searching in common directories
            common_dirs = ['.', './images', './assets', './resources']
            image_name = os.path.basename(image_path)

            for directory in common_dirs:
                potential_path = os.path.join(directory, image_name)
                if os.path.exists(potential_path):
                    image_path = potential_path
                    break
            else:
                return {
                    "error": f"Image file not found: {image_path}. Searched in {', '.join(common_dirs)}"
                }

        # Define multiple approaches to add the image
        def add_with_size():
            return ppt_utils.add_image(slide, image_path, left, top, width, height)

        def add_without_size():
            return ppt_utils.add_image(slide, image_path, left, top)

        def add_with_pil():
            from PIL import Image
            img = Image.open(image_path)
            img_width, img_height = img.size

            # Calculate aspect ratio and use it to determine missing dimension
            aspect_ratio = img_width / img_height

            if width is not None and height is None:
                h = width / aspect_ratio
                return ppt_utils.add_image(slide, image_path, left, top, width, h)
            elif height is not None and width is None:
                w = height * aspect_ratio
                return ppt_utils.add_image(slide, image_path, left, top, w, height)
            else:
                return ppt_utils.add_image(slide, image_path, left, top, width, height)

        approaches = [
            (add_with_size, "Adding image with specified dimensions"),
            (add_without_size, "Adding image with original dimensions"),
            (add_with_pil, "Adding image with calculated dimensions using PIL")
        ]

        picture, error = ppt_utils.try_multiple_approaches("add image", approaches)

        if error:
            return {
                "error": error
            }

        return {
            "message": f"Added image to slide {slide_index}",
            "shape_index": len(slide.shapes) - 1,
            "width": picture.width.inches,
            "height": picture.height.inches
        }

    def add_image_from_base64(
            self,
            slide_index: int,
            base64_string: str,
            left: float,
            top: float,
            width: Optional[float] = None,
            height: Optional[float] = None,
            presentation_id: Optional[str] = None
    ) -> Dict:
        """Add an image from a base64 encoded string to a slide."""
        # Validate
        err = self._validate()
        if err:
            return err

        pres = self._presentation

        # Check if slide index is valid
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            # Add the image
            picture = ppt_utils.add_image_from_base64(slide, base64_string, left, top, width, height)

            return {
                "message": f"Added image to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1,
                "width": picture.width.inches,
                "height": picture.height.inches
            }
        except Exception as e:
            return {
                "error": f"Failed to add image: {str(e)}"
            }

    def _export_slide_svg(self) -> str:
        """
        Export slide 0 as an SVG on disk (so Aspose writes it directly),
        then read it back as a string.
        """
        svg_path = None
        pptx_file = None
        # make sure Aspose knows about your system fonts
        slides.FontsLoader.load_external_fonts([
            "/Library/Fonts",
            "/System/Library/Fonts",
            os.path.expanduser("~/Library/Fonts")
        ])

        # choose your SVG options
        opts = SVGOptions.wysiwyg
        opts.vectorize_text = False
        opts.use_frame_size = True
        opts.metafile_rasterization_dpi = 300
        opts.external_fonts_handling = SvgExternalFontsHandling.EMBED

        # 1) save the in‐memory PPTX to disk
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_ppt:
            ppt_utils.save_presentation(self._presentation, tmp_ppt.name)
            pptx_file = tmp_ppt.name

        try:
            # 2) export slide 0 to a temp .svg file
            svg_fd, svg_path = tempfile.mkstemp(suffix=".svg")
            os.close(svg_fd)
            with open(svg_path, "wb") as f:
                with slides.Presentation(pptx_file) as as_pres:
                    # write_as_svg takes a file‐like or path
                    as_pres.slides[0].write_as_svg(f, opts)

            # 3) read it back
            with open(svg_path, "r", encoding="utf-8") as f:
                raw_svg = f.read()
        finally:
            if pptx_file: os.remove(pptx_file)
            if svg_path: os.remove(svg_path)

        return raw_svg

    def get_slide_image(self) -> bytes | dict:
        """
        1) Save current python-pptx presentation to a temp PPTX
        2) Export → SVG, strip watermark, convert → PNG
        3) Return raw PNG bytes
        """
        svg = self._export_slide_svg()
        clean_svg = remove_aspose_watermark(svg)
        png = svg_to_png(clean_svg)

        return png

    # ---- Table Tools ----
    def add_table(
            self,
            slide_index: int,
            rows: int,
            cols: int,
            left: float,
            top: float,
            width: float,
            height: float,
            data: Optional[List[List[str]]] = None,
    ) -> Dict:
        """Add a table to a slide with comprehensive parameter validation."""

        # Validate
        err = self._validate()
        if err:
            return err

        pres = self._presentation

        # Validate parameters
        valid, error = validate_parameters({
            "rows": (rows, [(is_positive, "must be a positive integer")]),
            "cols": (cols, [(is_positive, "must be a positive integer")]),
            "left": (left, [(is_non_negative, "must be non-negative")]),
            "top": (top, [(is_non_negative, "must be non-negative")]),
            "width": (width, [(is_positive, "must be positive")]),
            "height": (height, [(is_positive, "must be positive")]),
        })

        if not valid:
            return {"error": error}

        # Check if slide index is valid
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        # Validate data if provided
        if data is not None:
            if not isinstance(data, list):
                return {"error": "Data must be a list of rows"}

            for i, row in enumerate(data):
                if not isinstance(row, list):
                    return {"error": f"Row {i} must be a list of cell values"}

        try:
            # Add the table
            table, error = ppt_utils.safe_operation(
                "add_table",
                lambda: ppt_utils.add_table(slide, rows, cols, left, top, width, height)
            )

            if error:
                return {"error": error}

            # Populate the table if data is provided
            warnings = []
            if data:
                for row_idx, row_data in enumerate(data):
                    if row_idx >= rows:
                        warnings.append(
                            f"Ignored excess data: table has only {rows} rows but data has {len(data)} rows")
                        break

                    for col_idx, cell_text in enumerate(row_data):
                        if col_idx >= cols:
                            warnings.append(f"Ignored excess data in row {row_idx}: table has only {cols} columns")
                            break

                        _, cell_error = ppt_utils.safe_operation(
                            f"set_cell_text(row={row_idx}, col={col_idx})",
                            lambda: ppt_utils.set_cell_text(table, row_idx, col_idx, str(cell_text))
                        )

                        if cell_error:
                            warnings.append(cell_error)

            result = {
                "message": f"Added {rows}x{cols} table to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1
            }

            if warnings:
                result["warnings"] = warnings

            return result
        except Exception as e:
            return {
                "error": f"Failed to add table: {str(e)}"
            }

    def format_table_cell(
            self,
            slide_index: int,
            shape_index: int,
            row: int,
            col: int,
            font_size: Optional[int] = None,
            font_name: Optional[str] = None,
            bold: Optional[bool] = None,
            italic: Optional[bool] = None,
            color: Optional[List[int]] = None,
            bg_color: Optional[List[int]] = None,
            alignment: Optional[str] = None,
            vertical_alignment: Optional[str] = None,
    ) -> Dict:
        """Format a table cell with comprehensive error handling and parameter validation.

        This function applies formatting to a cell in a table on a slide. It provides options
        for text formatting (font size, name, style, color), cell background color, and
        text alignment.

        Args:
            slide_index: Index of the slide containing the table (0-based)
            shape_index: Index of the table shape on the slide (0-based)
            row: Row index of the cell to format (0-based)
            col: Column index of the cell to format (0-based)
            font_size: Font size in points
            font_name: Font name/family (e.g., 'Arial', 'Calibri')
            bold: Whether text should be bold (True/False)
            italic: Whether text should be italic (True/False)
            color: RGB color list for text [R, G, B] (0-255 for each value)
            bg_color: RGB color list for cell background [R, G, B] (0-255 for each value)
            alignment: Text alignment ('left', 'center', 'right', 'justify')
            vertical_alignment: Vertical text alignment ('top', 'middle', 'bottom')

        Returns:
            Dict with keys:
                - message: Success message
                - error: Error message if operation failed
                - warning: Warning message if some formatting was applied but some failed

        Examples:
            To format a header cell with bold text and gray background:
                format_table_cell(0, 1, 0, 1, font_size=14, bold=True, bg_color=[200, 200, 200])

            To center text in a cell:
                format_table_cell(0, 1, 2, 3, alignment="center", vertical_alignment="middle")
        """
        # Use the specified presentation or the current one
        err = self._validate()
        if err:
            return err

        pres = self._presentation

        # Check if slide index is valid
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        # Check if shape index is valid
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {
                "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
            }

        shape = slide.shapes[shape_index]

        # Validate parameters
        valid_alignments = ['left', 'center', 'right', 'justify']
        valid_vertical_alignments = ['top', 'middle', 'bottom']

        validations = {}

        if font_size is not None:
            validations["font_size"] = (font_size, [(is_positive, "must be a positive integer")])

        if alignment is not None:
            validations["alignment"] = (alignment.lower(), [(lambda x: x in valid_alignments,
                                                             f"must be one of {', '.join(valid_alignments)}")])

        if vertical_alignment is not None:
            validations["vertical_alignment"] = (vertical_alignment.lower(),
                                                 [(lambda x: x in valid_vertical_alignments,
                                                   f"must be one of {', '.join(valid_vertical_alignments)}")])

        if color is not None:
            validations["color"] = (color, [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")])

        if bg_color is not None:
            validations["bg_color"] = (bg_color,
                                       [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")])

        if validations:
            valid, error = validate_parameters(validations)
            if not valid:
                return {"error": error}

        try:
            # Check if shape is a table
            if not hasattr(shape, 'table'):
                # Try to recover if this is a graphic frame containing a table
                if hasattr(shape, 'graphic') and hasattr(shape.graphic, 'graphicData'):
                    # This might be a table in a graphic frame
                    warnings = ["Shape is not directly a table, attempting to extract table from graphic frame"]
                    # Further recovery logic would be needed here
                    return {
                        "error": "Shape at index {shape_index} is not a table",
                        "warning": "If this is a table, it might be in a graphic frame which requires different handling"
                    }
                else:
                    return {
                        "error": f"Shape at index {shape_index} is not a table"
                    }

            table = shape.table

            # Check if row and column indices are valid
            if row < 0 or row >= len(table.rows):
                return {
                    "error": f"Invalid row index: {row}. Available rows: 0-{len(table.rows) - 1}"
                }

            if col < 0 or col >= len(table.columns):
                return {
                    "error": f"Invalid column index: {col}. Available columns: 0-{len(table.columns) - 1}"
                }

            # Get the cell
            cell = table.cell(row, col)

            # Format the cell with error handling
            warnings = []

            # Try multiple formatting operations and collect any warnings
            try:
                ppt_utils.format_table_cell(
                    cell,
                    font_size=font_size,
                    font_name=font_name,
                    bold=bold,
                    italic=italic,
                    color=tuple(color) if color else None,
                    bg_color=tuple(bg_color) if bg_color else None,
                    alignment=alignment,
                    vertical_alignment=vertical_alignment
                )
            except Exception as e:
                # Try individual formatting operations to recover
                formatting_ops = [
                    (lambda: ppt_utils.format_text(cell.text_frame, font_size=font_size, font_name=font_name,
                                                   bold=bold, italic=italic,
                                                   color=tuple(color) if color else None,
                                                   alignment=alignment),
                     "text formatting"),

                    (lambda: cell.fill.solid() if bg_color else None, "background preparation"),

                    (lambda: setattr(cell.fill.fore_color, 'rgb',
                                     tuple(bg_color)) if bg_color else None,
                     "background color"),

                    (lambda: setattr(cell.text_frame, 'vertical_anchor',
                                     ppt_utils.vertical_alignment_map.get(vertical_alignment))
                    if vertical_alignment else None,
                     "vertical alignment")
                ]

                for op_func, op_name in formatting_ops:
                    try:
                        op_func()
                    except Exception as sub_e:
                        warnings.append(f"Failed to apply {op_name}: {str(sub_e)}")

            result = {
                "message": f"Formatted cell at row {row}, column {col} in table at shape index {shape_index} on slide {slide_index}"
            }

            if warnings:
                result["warnings"] = warnings

            return result
        except Exception as e:
            return {
                "error": f"Failed to format table cell: {str(e)}"
            }

    # ---- Shape Tools ----
    def add_shape(
            self,
            slide_index: int,
            shape_type: str,
            left: float,
            top: float,
            width: float,
            height: float,
            fill_color: Optional[List[int]] = None,
            line_color: Optional[List[int]] = None,
            line_width: Optional[float] = None
    ) -> Dict:
        """Add an auto shape to a slide.

        This function adds a shape to a slide in the presentation. It supports various shape types
        and allows customization of fill color, line color, and line width.

        Args:
            slide_index: Index of the slide to add the shape to (0-based)
            shape_type: Type of shape to add. Supported types include:
                - Basic shapes: 'rectangle', 'rounded_rectangle', 'oval', 'triangle', 'diamond'
                - Polygons: 'pentagon', 'hexagon', 'heptagon', 'octagon'
                - Stars and arrows: 'star', 'arrow'
                - Misc: 'cloud', 'heart', 'lightning_bolt', 'sun', 'moon', 'smiley_face', 'no_symbol'
                - Flowchart: 'flowchart_process', 'flowchart_decision', 'flowchart_data'
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            fill_color: RGB color list for shape fill [R, G, B] (0-255 for each value)
            line_color: RGB color list for shape outline [R, G, B] (0-255 for each value)
            line_width: Width of the shape outline in points
        """
        # Use the specified presentation or the current one
        # Validate presentation
        err = self._validate()
        if err:
            return err

        pres = self._presentation

        # Check if slide index is valid
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            # Use the direct implementation that bypasses the enum issues
            shape = add_shape_direct(slide, shape_type, left, top, width, height)

            # Format the shape if formatting options are provided
            if any([fill_color, line_color, line_width]):
                ppt_utils.format_shape(
                    shape,
                    fill_color=tuple(fill_color) if fill_color else None,
                    line_color=tuple(line_color) if line_color else None,
                    line_width=line_width
                )

            return {
                "message": f"Added {shape_type} shape to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1
            }
        except ValueError as e:
            # Specific handling for validation errors
            return {
                "error": str(e)
            }
        except Exception as e:
            return {
                "error": f"Failed to add shape '{shape_type}': {str(e)}"
            }

    # ---- Chart Tools ----
    def add_chart(
            self,
            slide_index: int,
            chart_type: str,
            left: float,
            top: float,
            width: float,
            height: float,
            categories: List[str],
            series_names: List[str],
            series_values: List[List[float]],
            has_legend: bool = True,
            legend_position: str = "right",
            has_data_labels: bool = False,
            title: Optional[str] = None,
    ) -> Dict:
        """Add a chart to a slide with comprehensive error handling."""
        # Validate presentation
        err = self._validate()
        if err:
            return err

        pres = self._presentation

        # Check if slide index is valid
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        # Validate chart type
        valid_chart_types = [
            'column', 'stacked_column', 'bar', 'stacked_bar', 'line',
            'line_markers', 'pie', 'doughnut', 'area', 'stacked_area',
            'scatter', 'radar', 'radar_markers'
        ]
        if chart_type.lower() not in valid_chart_types:
            return {
                "error": f"Invalid chart type: '{chart_type}'. Valid types are: {', '.join(valid_chart_types)}"
            }

        # Validate series data
        if len(series_names) != len(series_values):
            return {
                "error": f"Number of series names ({len(series_names)}) must match number of series values ({len(series_values)})"
            }

        # Validate categories list
        if not categories:
            return {
                "error": "Categories list cannot be empty"
            }

        # Validate that all series have the same number of values as categories
        for i, values in enumerate(series_values):
            if len(values) != len(categories):
                return {
                    "error": f"Series '{series_names[i]}' has {len(values)} values but there are {len(categories)} categories"
                }

        try:
            # Add the chart
            chart, error = ppt_utils.safe_operation(
                "add_chart",
                lambda: ppt_utils.add_chart(
                    slide, chart_type, left, top, width, height,
                    categories, series_names, series_values
                )
            )

            if error:
                return {"error": error}

            # Format the chart
            _, error = ppt_utils.safe_operation(
                "format_chart",
                lambda: ppt_utils.format_chart(
                    chart,
                    has_legend=has_legend,
                    legend_position=legend_position,
                    has_data_labels=has_data_labels,
                    title=title
                )
            )

            if error:
                return {
                    "warning": f"Chart created but failed to format: {error}",
                    "shape_index": len(slide.shapes) - 1
                }

            return {
                "message": f"Added {chart_type} chart to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1
            }
        except Exception as e:
            return {
                "error": f"Failed to add chart: {str(e)}"
            }

    def move_element(self, shape_index: int, left: float, top: float) -> Dict:
        """Tool to reposition an existing shape on slide 0."""
        # validate slide exists
        err = self._validate(slide_index=0)
        if err:
            return err

        slide = self._presentation.slides[0]
        try:
            moved = ppt_utils.move_shape(slide, shape_index, left, top)
            # return new coords so user can verify
            return {
                "message": f"Moved shape {shape_index} → ({left}\", {top}\")",
                "new_left": moved.left.inches,
                "new_top":  moved.top.inches
            }
        except Exception as e:
            return {"error": str(e)}